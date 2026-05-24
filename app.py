# ============================================================================
# AI视频生成应用 - 主程序文件
# ============================================================================
# 功能说明：
# 1. 支持多种文档格式上传（Word、PDF、PPT、TXT、音频等）
# 2. 使用AI生成专业的视频prompt
# 3. 调用Seedance API生成真实视频
# 4. 本地降级方案：使用OpenCV生成演示视频
# ============================================================================

# 标准库导入
import os
import json
import time
from datetime import datetime
import tempfile

# 第三方库导入
from flask import Flask, request, jsonify, send_file, send_from_directory
from flask_cors import CORS
import requests
from dotenv import load_dotenv

# 自定义模块导入
from logger_config import setup_logger

# 加载环境变量
load_dotenv()

# ============================================================================
# 日志系统初始化
# ============================================================================
# 日志文件位置：logs/app_YYYYMMDD.log（所有日志）
#              logs/error_YYYYMMDD.log（仅错误日志）
# 日志格式：[时间] [级别] [模块] 消息
# ============================================================================
logger = setup_logger('app')
logger.info("="*60)
logger.info("应用启动")
logger.info("="*60)

# ============================================================================
# Flask应用初始化
# ============================================================================
app = Flask(__name__)
# 启用跨域资源共享，允许前端访问API
CORS(app)

# ============================================================================
# 应用配置
# ============================================================================

# 配置文件上传大小限制（500MB）
# 注意：修改此值时，需要同时考虑服务器内存和磁盘空间
app.config['MAX_CONTENT_LENGTH'] = 500 * 1024 * 1024

# 创建必要的目录
# UPLOAD_FOLDER: 存储上传的原始文件
# OUTPUT_FOLDER: 存储生成的视频文件
UPLOAD_FOLDER = 'uploads'
OUTPUT_FOLDER = 'outputs'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

# ============================================================================
# AI视频生成模型配置
# ============================================================================

# 支持的环境变量：
# SEEDANCE_API_KEY: 字节跳动API密钥
# SEEDANCE_ENDPOINT: 字节跳动API端点
# RUNWAY_API_KEY: RunwayML API密钥
SEEDANCE_API_KEY = os.environ.get('SEEDANCE_API_KEY') or 'your_seedance_api_key'
SEEDANCE_ENDPOINT = os.environ.get('SEEDANCE_ENDPOINT') or 'your_seedance_endpoint'

# AI视频生成模型配置字典
# 每个模型包含：
# - name: 模型显示名称
# - api_url: 文本生成API地址
# - video_api_url: 视频生成API地址
# - api_key: API密钥
# - model: 模型标识符
# - video_model: 视频生成模型名称
# - description: 模型描述
AI_VIDEO_MODELS = {
    'seedance': {
        'name': 'Seedance 1.5 Pro (字节跳动)',
        'api_url': 'https://ark.cn-beijing.volces.com/api/v3/chat/completions',
        'video_api_url': 'https://ark.cn-beijing.volces.com/api/v3/contents/generations/tasks',
        'api_key': SEEDANCE_API_KEY,
        'model': SEEDANCE_ENDPOINT,
        'video_model': 'seedance-1.5-pro',
        'description': '字节跳动 Seedance 1.5 Pro AI视频生成引擎'
    },
    'runway': {
        'name': 'RunwayML Gen-2 (国内版)',
        'api_url': 'https://api.runwayml.com/v1/generate',
        'api_key': os.environ.get('RUNWAY_API_KEY', 'YOUR_RUNWAY_API_KEY'),
        'model': 'gen2',
        'description': '国内可访问的Runway视频生成模型'
    },
    'pika': {
        'name': 'Pika Labs (国内版)',
        'api_url': 'https://api.pika.art/v1/generate',
        'api_key': os.environ.get('PIKA_API_KEY', 'YOUR_PIKA_API_KEY'),
        'model': 'pika-1.0',
        'description': '国内可访问的Pika视频生成模型'
    }
}

# ============================================================================
# Flask路由定义
# ============================================================================

@app.route('/')
def index():
    """
    主页路由
    返回index.html文件
    """
    return send_from_directory('.', 'index.html')

@app.route('/styles.css')
def styles():
    """
    CSS样式文件路由
    返回styles.css文件
    """
    return send_from_directory('.', 'styles.css')

@app.route('/script.js')
def script():
    """
    JavaScript文件路由
    返回script.js文件
    """
    return send_from_directory('.', 'script.js')

# ============================================================================
# API路由
# ============================================================================

@app.route('/api/upload', methods=['POST'])
def upload_file():
    """
    文件上传API

    功能：
    1. 验证文件大小（最大500MB）
    2. 验证文件类型（支持PDF、PPT、Word、TXT、音频等）
    3. 保存文件到uploads目录
    4. 返回文件信息

    返回JSON格式：
    {
        'message': '文件上传成功',
        'filename': '20240410143000_document.pdf',
        'filepath': 'uploads/20240410143000_document.pdf',
        'fileSize': 1234567,
        'fileType': 'document' 或 'audio'
    }
    """
    if request.content_length > app.config['MAX_CONTENT_LENGTH']:
        return jsonify({'error': f'文件大小超过限制（最大{app.config["MAX_CONTENT_LENGTH"]//1024//1024}MB）'}), 413

    if 'file' not in request.files:
        return jsonify({'error': '没有文件部分'}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': '没有选择文件'}), 400

    # 检查文件类型
    allowed_extensions = {'.pdf', '.ppt', '.pptx', '.doc', '.docx', '.txt', '.wav', '.mp3', '.ogg', '.webm'}
    file_ext = os.path.splitext(file.filename)[1].lower()

    if file_ext not in allowed_extensions:
        return jsonify({'error': f'不支持的文件格式: {file_ext}'}), 400

    # 保存上传的文件
    timestamp = datetime.now().strftime('%Y%m%d%H%M%S')
    filename = f"{timestamp}_{file.filename}"
    filepath = os.path.join(UPLOAD_FOLDER, filename)
    file.save(filepath)

    return jsonify({
        'message': '文件上传成功',
        'filename': filename,
        'filepath': filepath,
        'fileSize': os.path.getsize(filepath),
        'fileType': 'audio' if file_ext in {'.wav', '.mp3', '.ogg', '.webm'} else 'document'
    })

@app.route('/api/generate-video', methods=['POST'])
def generate_video():
    """
    视频生成API

    功能：
    1. 接收前端传递的参数（文件名、AI模型、风格、时长、旁白）
    2. 读取上传的文件内容
    3. 使用AI生成专业的视频prompt
    4. 调用视频生成API或本地生成引擎
    5. 返回视频URL

    请求参数：
    {
        'filename': '上传的文件名',
        'aiModel': 'seedance' 或 'runway' 或 'pika',
        'videoStyle': 'business' 或 'creative' 或 'minimalist' 或 'tech',
        'videoDuration': 'short' 或 'medium' 或 'long',
        'narrator': '旁白声音类型',
        'fileType': 'document' 或 'audio'
    }

    返回JSON格式：
    {
        'message': '视频生成成功',
        'videoFilename': '20240410143000_document_seedance_1234567890.mp4',
        'videoUrl': '/api/video/20240410143000_document_seedance_1234567890.mp4',
        'prompt': '生成的视频prompt'
    }
    """
    data = request.json
    print(f"\n=== 开始生成视频请求 ===")
    print(f"请求数据: {data}")

    # 获取参数
    ai_model = data.get('aiModel', 'seedance')
    video_style = data.get('videoStyle', 'business')
    video_duration = data.get('videoDuration', 'medium')
    narrator = data.get('narrator', 'female1')
    filename = data.get('filename')
    file_type = data.get('fileType', 'document')

    if not filename:
        return jsonify({'error': '缺少文件名'}), 400

    filepath = os.path.join(UPLOAD_FOLDER, filename)
    if not os.path.exists(filepath):
        return jsonify({'error': '文件不存在'}), 400

    try:
        # 处理原始素材内容
        content = ""
        file_ext = os.path.splitext(filename)[1].lower()
        file_size = os.path.getsize(filepath)
        print(f"文件类型: {file_ext}")
        print(f"文件大小: {file_size} bytes")

        MAX_CONTENT_LENGTH = 100000

        if file_size > 50 * 1024 * 1024:
            print("[WARN] 文件较大，将只读取部分内容")

        if file_type == 'audio' or file_ext in ['.wav', '.mp3', '.ogg', '.webm']:
            # 音频文件：语音转文本
            print("[INFO] 检测到音频文件，开始语音转文本...")
            content = transcribe_audio(filepath, filename, file_size)
            print(f"[OK] 语音转文本完成，长度: {len(content)}")
        elif file_ext == '.txt':
            with open(filepath, 'r', encoding='utf-8') as f:
                content = f.read(MAX_CONTENT_LENGTH)
                if len(content) == MAX_CONTENT_LENGTH:
                    print(f"[WARN] 文件内容过长，已截取前{MAX_CONTENT_LENGTH}字符")
        elif file_ext in ['.doc', '.docx']:
            try:
                from docx import Document
                doc = Document(filepath)
                paragraphs = [paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()]
                max_paragraphs = 500
                if len(paragraphs) > max_paragraphs:
                    print(f"[WARN] Word文档段落数过多，已截取前{max_paragraphs}段")
                    paragraphs = paragraphs[:max_paragraphs]
                content = '\n'.join(paragraphs)
            except (ImportError, ValueError, Exception) as e:
                if file_ext == '.doc':
                    print(f"[WARN] 旧的 .doc 格式（二进制）不支持，请使用 .docx 格式")
                    raise ValueError("不支持旧的 .doc 二进制格式，请转换为 .docx 格式后重试")

                try:
                    import zipfile
                    import xml.etree.ElementTree as ET
                    with zipfile.ZipFile(filepath) as z:
                        xml_content = z.read('word/document.xml')
                        tree = ET.fromstring(xml_content)
                        content = '\n'.join([node.text for node in tree.iter() if node.text])
                except Exception as zip_error:
                    print(f"[ERROR] 无法读取Word文档: {str(zip_error)}")
                    raise ValueError(f"无法读取Word文档，请确保文件格式正确: {str(e)}")
        elif file_ext == '.pdf':
            try:
                import pypdf
                with open(filepath, 'rb') as f:
                    pdf_reader = pypdf.PdfReader(f)
                    max_pages = 50
                    total_pages = len(pdf_reader.pages)
                    if total_pages > max_pages:
                        print(f"[WARN] PDF页数过多，已截取前{max_pages}页")
                        pages_to_read = pdf_reader.pages[:max_pages]
                    else:
                        pages_to_read = pdf_reader.pages
                    content = '\n'.join([page.extract_text() or '' for page in pages_to_read])
            except ImportError:
                raise ValueError("请安装pypdf库以支持PDF文件处理")
        elif file_ext == '.pptx':
            try:
                from pptx import Presentation
                presentation = Presentation(filepath)
                slides_text = []
                max_slides = 50
                total_slides = len(presentation.slides)
                if total_slides > max_slides:
                    print(f"[WARN] PPT幻灯片数过多，已截取前{max_slides}张")
                    slides_to_read = list(presentation.slides)[:max_slides]
                else:
                    slides_to_read = presentation.slides

                for slide in slides_to_read:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text.strip():
                            slides_text.append(shape.text.strip())
                content = '\n'.join(slides_text)
            except ImportError:
                raise ValueError("请安装python-pptx库以支持PPTX文件处理")
        else:
            raise ValueError(f"不支持的文件格式: {file_ext}")

        print(f"读取原始素材内容成功，长度: {len(content)}")

        if len(content) > MAX_CONTENT_LENGTH:
            print(f"[WARN] 内容过长，已截取至{MAX_CONTENT_LENGTH}字符")
            content = content[:MAX_CONTENT_LENGTH]

        # 生成专业的AI视频prompt
        print("\n[Stage 1] 生成专业的AI视频prompt...")
        video_prompt = generate_video_prompt(content, video_style, video_duration, narrator)
        print(f"[OK] Prompt生成完成，长度: {len(video_prompt)}")

        # 生成视频
        video_filename = f"{os.path.splitext(filename)[0]}_{ai_model}_{int(time.time())}.mp4"
        video_filepath = os.path.join(OUTPUT_FOLDER, video_filename)
        print(f"视频保存路径: {video_filepath}")

        # 调用AI视频生成模型
        model_info = AI_VIDEO_MODELS.get(ai_model)
        if not model_info:
            raise ValueError(f"不支持的AI模型: {ai_model}")

        print(f"使用模型: {model_info['name']}")

        # 调用视频生成API
        generate_video_with_model(model_info, content, video_prompt, video_style, video_duration, narrator, video_filepath)

        print(f"视频生成完成，检查文件: {os.path.exists(video_filepath)}")

        return jsonify({
            'message': '视频生成成功',
            'videoFilename': video_filename,
            'videoUrl': f'/api/video/{video_filename}',
            'prompt': video_prompt
        })
    except Exception as e:
        print(f"[ERROR] 生成视频失败: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500

@app.route('/api/video/<filename>', methods=['GET'])
def get_video(filename):
    """获取生成的视频"""
    filepath = os.path.join(OUTPUT_FOLDER, filename)
    if not os.path.exists(filepath):
        return jsonify({'error': '视频不存在'}), 404

    return send_file(filepath, as_attachment=True)

def transcribe_audio(filepath, filename, file_size):
    """音频转文本"""
    try:
        import speech_recognition as sr
        recognizer = sr.Recognizer()

        with sr.AudioFile(filepath) as source:
            audio = recognizer.record(source)
            text = recognizer.recognize_google(audio, language='zh-CN')
            print(f"[OK] 音频转录成功: {len(text)} 字符")
            return text
    except Exception as e:
        print(f"[WARN] 音频转录失败: {str(e)}")

    # 降级方案：返回基于文件信息的文本
    return f"""音频文件: {filename}
文件大小: {file_size / 1024 / 1024:.2f} MB
音频格式: {os.path.splitext(filename)[1].upper()}

注: 为了获得完整的转录内容，请确保已安装speech_recognition库和ffmpeg。
当前模式下，系统将基于视频样式和时长生成演示视频。
实际使用中，音频内容将被转录并整合到视频中。"""

def analyze_document_structure(content_lines):
    """智能分析文档结构"""
    structure = {
        'title': '文档介绍',
        'key_points': [],
        'summary': '',
        'sections': []
    }

    # 提取标题（第一行非空行）
    for line in content_lines:
        if line.strip() and len(line.strip()) > 5:
            structure['title'] = line.strip()
            break

    # 提取关键要点
    for line in content_lines:
        line = line.strip()
        if line and len(line) > 10:
            # 识别要点标记（如：•、-、*、数字等）
            if any(marker in line[:3] for marker in ['•', '-', '*', '1.', '2.', '3.']):
                structure['key_points'].append(line)
            elif len(structure['key_points']) < 8:  # 最多提取8个要点
                structure['key_points'].append(line)

    # 生成摘要（取前几个要点）
    if structure['key_points']:
        structure['summary'] = ' '.join(structure['key_points'][:3])

    return structure

def select_content_for_duration(key_points, num_scenes, words_per_scene):
    """根据视频时长选择合适的内容"""
    if not key_points:
        return []

    # 确保每个场景有合适的内容量
    selected = []
    points_per_scene = max(1, len(key_points) // num_scenes)

    for i in range(num_scenes):
        start_idx = i * points_per_scene
        end_idx = start_idx + points_per_scene

        if start_idx < len(key_points):
            scene_points = key_points[start_idx:end_idx]
            if scene_points:
                # 合并多个要点为一个场景
                combined = ' '.join(scene_points[:2])  # 每个场景最多2个要点
                selected.append(combined)

    # 确保至少有足够的内容
    while len(selected) < num_scenes and len(selected) < len(key_points):
        selected.append(key_points[len(selected)])

    return selected[:num_scenes]

def build_scene_descriptions(document_title, selected_points, style_desc, duration_info):
    """构建场景描述"""
    scenes = []
    scene_duration = duration_info['seconds'] // duration_info['scenes']

    # 场景1：开场 - 展示文档标题
    scenes.append(f"""场景1（开场，{scene_duration}秒）：
- 画面：文档标题"{document_title}"以大字体居中显示
- 背景：{style_desc.split('，')[0]}的渐变背景
- 动画：标题从下方淡入，配合光效
- 文字：标题文字清晰，使用{style_desc.split('，')[1] if '，' in style_desc else '专业'}字体
- 配音：介绍文档主题和目的（约{scene_duration * 3}字）""")

    # 场景2-N：展示核心内容
    for i, point in enumerate(selected_points, 2):
        if i > duration_info['scenes']:
            break

        scenes.append(f"""场景{i}（内容展示，{scene_duration}秒）：
- 画面：展示要点内容"{point[:30]}..."
- 布局：左侧文字，右侧配图或图标
- 动画：文字逐行显示，配合图标动画
- 颜色：使用{style_desc.split('，')[0] if '，' in style_desc else '专业'}配色
- 配音：详细说明要点内容（约{scene_duration * 3}字）""")

    # 最后场景：结尾 - 总结和行动号召
    scenes.append(f"""场景{len(scenes)+1}（结尾，{scene_duration}秒）：
- 画面：文档标题再次出现，下方显示"谢谢观看"
- 背景：与开场呼应的渐变效果
- 动画：标题和文字同时淡入
- 文字：简洁明了，突出重点
- 配音：总结文档价值，感谢观看（约{scene_duration * 3}字）""")

    return '\n\n'.join(scenes)

def generate_video_prompt(content, video_style, video_duration, narrator):
    """根据文档内容生成智能的文档介绍视频prompt"""
    print(f"\n=== generate_video_prompt 开始 ===")
    print(f"视频风格: {video_style}")
    print(f"视频时长: {video_duration}")
    print(f"配音风格: {narrator}")
    print(f"内容长度: {len(content)} 字符")

    # 风格映射 - 专注于文档介绍视频
    style_descriptions = {
        'business': '专业商务风格，简洁大方，适合企业文档介绍，使用蓝色和灰色调，清晰的排版',
        'creative': '创意艺术风格，充满活力，适合创意文档，使用鲜艳色彩和动态效果',
        'minimalist': '极简风格，干净整洁，突出核心信息，使用白色背景和黑色文字，留白充足',
        'tech': '科技风格，现代感强，适合技术文档，使用深色背景和霓虹色点缀'
    }

    # 时长映射 - 针对文档介绍视频优化
    duration_maps = {
        'short': {'seconds': 30, 'scenes': 2, 'desc': '30秒短视频', 'words_per_scene': 15},
        'medium': {'seconds': 60, 'scenes': 3, 'desc': '60秒中等视频', 'words_per_scene': 20},
        'long': {'seconds': 100, 'scenes': 4, 'desc': '100秒长视频', 'words_per_scene': 25}
    }

    # 配音映射
    narrator_descriptions = {
        'female1': '温柔专业的女声，语速适中，适合商务文档介绍',
        'female2': '活泼亲切的女声，语速稍快，适合创意内容',
        'male1': '沉稳权威的男声，语速稳健，适合技术文档',
        'male2': '温和友好的男声，语速适中，适合通用场景'
    }

    style_desc = style_descriptions.get(video_style, style_descriptions['business'])
    duration_info = duration_maps.get(video_duration, duration_maps['medium'])
    narrator_desc = narrator_descriptions.get(narrator, narrator_descriptions['female1'])

    # 智能分析文档内容
    content_lines = content.strip().split('\n')

    # 提取文档结构
    document_structure = analyze_document_structure(content_lines)

    # 提取标题
    document_title = document_structure.get('title', '文档介绍')

    # 提取核心要点
    key_points = document_structure.get('key_points', [])

    # 提取摘要
    document_summary = document_structure.get('summary', '')

    # 根据时长选择合适的内容量
    selected_points = select_content_for_duration(key_points, duration_info['scenes'], duration_info['words_per_scene'])

    # 构建场景描述
    scene_descriptions = build_scene_descriptions(
        document_title,
        selected_points,
        style_desc,
        duration_info
    )

    # 生成最终的prompt
    video_prompt = f"""【文档介绍视频生成】

视频主题：{document_title}

【场景设计】
{scene_descriptions}

【视觉风格】
{style_desc}

【配音要求】
{narrator_desc}

【技术规格】
- 总时长：{duration_info['desc']}
- 场景数量：{duration_info['scenes']}个
- 每个场景约{duration_info['seconds'] // duration_info['scenes']}秒
- 文字叠加：清晰易读，突出关键信息
- 转场效果：平滑自然的场景切换
- 背景音乐：轻柔的背景音，不干扰配音

【内容要求】
- 简洁明了，突出文档核心价值
- 每个场景聚焦一个主要观点
- 使用图表、图标等视觉元素辅助说明
- 保持一致的视觉风格和配色
- 确保信息传达清晰准确"""

    print(f"[OK] Prompt生成完成")
    return video_prompt

def generate_video_with_model(model_info, content, video_prompt, video_style, video_duration, narrator, video_filepath):
    """
    使用AI视频生成模型生成视频

    参数：
        model_info: AI模型配置字典
        content: 文档原始内容
        video_prompt: AI生成的视频prompt
        video_style: 视频风格（business/creative/minimalist/tech）
        video_duration: 视频时长（short/medium/long）
        narrator: 旁白类型
        video_filepath: 视频保存路径

    功能：
        1. 根据模型类型选择生成方式
        2. Seedance模型：调用API生成真实视频
        3. 其他模型：使用本地生成引擎
        4. 支持降级机制
    """
    print(f"\n=== 开始生成视频 ===")
    print(f"使用模型: {model_info['name']}")
    print(f"文档内容长度: {len(content)} 字符")
    print(f"生成的Prompt长度: {len(video_prompt)} 字符")

    # 将文档内容和生成的Prompt组合作为视频内容
    video_data = {
        'original_content': content,
        'ai_prompt': video_prompt,
        'style': video_style,
        'duration': video_duration,
        'narrator': narrator
    }

    # 调用 Seedance 1.5-pro 生成真实视频
    if model_info.get('model') == SEEDANCE_ENDPOINT:
        print(f"调用 Seedance 1.5-pro API 生成视频...")
        generate_video_with_seedance_api(model_info, video_prompt, video_data, video_filepath)
    else:
        print("使用本地视频生成引擎...")
        create_intelligent_video(video_filepath, video_data)

def generate_video_with_seedance_api(model_info, video_prompt, video_data, video_filepath):
    """
    使用 Seedance 1.5-pro API 生成真实视频 - 增强错误处理

    参数：
        model_info: AI模型配置字典
        video_prompt: AI生成的视频prompt
        video_data: 视频数据字典（包含内容、风格、时长等）
        video_filepath: 视频保存路径

    功能：
        1. 验证API配置（密钥、端点）
        2. 构建API请求参数
        3. 调用Seedance API创建视频生成任务
        4. 轮询查询任务状态
        5. 下载生成的视频
        6. 支持重试机制（创建任务、查询状态、下载视频）
        7. 支持降级机制（API失败时使用本地生成）

    重试策略：
        - 创建任务：最多3次，指数退避（5-30秒）
        - 查询状态：最多30次，间隔10秒
        - 下载视频：最多3次，间隔5秒

    错误处理：
        - 429错误：等待更长时间后重试
        - 5xx错误：自动重试
        - 超时错误：自动重试
        - 连接错误：自动重试
    """
    api_key = model_info.get('api_key')
    video_api_url = model_info.get('video_api_url')
    endpoint = model_info.get('model')

    # 验证API配置
    if not api_key or not endpoint or not video_api_url:
        logger.warning("Seedance API 配置不完整:")
        if not api_key:
            logger.warning("  - 缺少API密钥")
        if not endpoint:
            logger.warning("  - 缺少端点")
        if not video_api_url:
            logger.warning("  - 缺少视频API URL")
        logger.info("降级到本地视频生成...")
        create_intelligent_video(video_filepath, video_data)
        return

    logger.info(f"Seedance 1.5-pro API 端点: {endpoint}")
    logger.info(f"视频API URL: {video_api_url}")
    logger.info(f"生成的视频Prompt:\n{video_prompt[:500]}...")

    # 从文档内容中提取关键信息作为视频描述
    original_content = video_data.get('original_content', '')
    content_lines = [line.strip() for line in original_content.split('\n') if line.strip() and len(line.strip()) > 10]
    key_content = '\n'.join(content_lines[:5])

    # 构建视频生成参数
    duration = video_data.get('duration', 'medium')
    duration_map = {'short': 5, 'medium': 8, 'long': 12}
    video_duration = duration_map.get(duration, 8)

    # 构建图生视频的prompt
    style = video_data.get('style', 'business')
    style_prompts = {
        'business': '专业商务风格，简洁大方，适合企业宣传',
        'creative': '创意艺术风格，充满想象力，视觉冲击力强',
        'minimalist': '极简风格，干净简约，突出核心信息',
        'tech': '科技感风格，现代科技元素，未来感设计'
    }
    style_desc = style_prompts.get(style, style_prompts['business'])

    # 组合最终的文本prompt
    final_prompt = f"{key_content}\n{style_desc}\n{video_prompt}"

    # 限制prompt长度
    if len(final_prompt) > 500:
        final_prompt = final_prompt[:500]
        logger.info(f"Prompt已截断至500字符")

    # 构建请求payload
    payload = {
        'model': endpoint,
        'content': [
            {
                'type': 'text',
                'text': f"{final_prompt} --duration {video_duration} --camerafixed false --watermark true"
            }
        ]
    }

    headers = {
        'Content-Type': 'application/json',
        'Authorization': f'Bearer {api_key}'
    }

    # 重试配置
    max_create_retries = 3
    base_retry_delay = 5  # 初始重试延迟（秒）
    max_retry_delay = 30   # 最大重试延迟（秒）

    # 创建任务（带重试）
    task_id = None
    for retry in range(max_create_retries):
        try:
            logger.info(f"调用Seedance图生视频API创建任务 (尝试 {retry + 1}/{max_create_retries})...")

            # 设置超时时间：连接超时10秒，读取超时60秒
            response = requests.post(
                video_api_url, 
                json=payload, 
                headers=headers, 
                timeout=(10, 60)
            )

            if response.status_code != 200:
                error_msg = f"API返回状态码: {response.status_code}"
                try:
                    error_detail = response.json()
                    error_msg += f", 错误详情: {error_detail}"
                except:
                    error_msg += f", 响应内容: {response.text[:200]}"

                logger.warning(f"Seedance API 调用失败: {error_msg}")

                # 如果是429（请求过多），等待更长时间
                if response.status_code == 429:
                    retry_delay = min(base_retry_delay * (2 ** retry) * 2, max_retry_delay)
                    logger.info(f"遇到请求限制，等待 {retry_delay} 秒后重试...")
                    time.sleep(retry_delay)
                    continue

                # 如果是5xx错误，可以重试
                if 500 <= response.status_code < 600:
                    if retry < max_create_retries - 1:
                        retry_delay = min(base_retry_delay * (2 ** retry), max_retry_delay)
                        logger.info(f"服务器错误，等待 {retry_delay} 秒后重试...")
                        time.sleep(retry_delay)
                        continue

                # 其他错误不再重试
                logger.info("降级到本地视频生成...")
                create_intelligent_video(video_filepath, video_data)
                return

            # 解析响应
            try:
                result = response.json()
                task_id = result.get('id')

                if not task_id:
                    logger.warning("Seedance API 未返回任务ID")
                    logger.info(f"响应内容: {result}")
                    if retry < max_create_retries - 1:
                        retry_delay = min(base_retry_delay * (2 ** retry), max_retry_delay)
                        logger.info(f"等待 {retry_delay} 秒后重试...")
                        time.sleep(retry_delay)
                        continue
                    else:
                        logger.info("降级到本地视频生成...")
                        create_intelligent_video(video_filepath, video_data)
                        return

                logger.info(f"任务创建成功，任务ID: {task_id}")
                break

            except json.JSONDecodeError as e:
                logger.warning(f"解析响应JSON失败: {str(e)}")
                logger.info(f"响应内容: {response.text[:200]}")
                if retry < max_create_retries - 1:
                    retry_delay = min(base_retry_delay * (2 ** retry), max_retry_delay)
                    logger.info(f"等待 {retry_delay} 秒后重试...")
                    time.sleep(retry_delay)
                    continue
                else:
                    logger.info("降级到本地视频生成...")
                    create_intelligent_video(video_filepath, video_data)
                    return

        except requests.exceptions.Timeout as e:
            logger.warning(f"API请求超时: {str(e)}")
            if retry < max_create_retries - 1:
                retry_delay = min(base_retry_delay * (2 ** retry), max_retry_delay)
                logger.info(f"等待 {retry_delay} 秒后重试...")
                time.sleep(retry_delay)
                continue
            else:
                logger.info("降级到本地视频生成...")
                create_intelligent_video(video_filepath, video_data)
                return

        except requests.exceptions.ConnectionError as e:
            logger.warning(f"API连接错误: {str(e)}")
            if retry < max_create_retries - 1:
                retry_delay = min(base_retry_delay * (2 ** retry), max_retry_delay)
                logger.info(f"等待 {retry_delay} 秒后重试...")
                time.sleep(retry_delay)
                continue
            else:
                logger.info("降级到本地视频生成...")
                create_intelligent_video(video_filepath, video_data)
                return

        except Exception as e:
            logger.error(f"创建任务时发生未预期的错误: {str(e)}")
            import traceback
            traceback.print_exc()
            if retry < max_create_retries - 1:
                retry_delay = min(base_retry_delay * (2 ** retry), max_retry_delay)
                logger.info(f"等待 {retry_delay} 秒后重试...")
                time.sleep(retry_delay)
                continue
            else:
                logger.info("降级到本地视频生成...")
                create_intelligent_video(video_filepath, video_data)
                return

    # 轮询查询任务状态（带重试）
    if task_id:
        max_poll_retries = 30
        poll_interval = 10
        poll_timeout = 30  # 每次查询的超时时间

        print(f"[INFO] 开始轮询任务状态，最大等待时间: {max_poll_retries * poll_interval} 秒")

        for attempt in range(max_poll_retries):
            try:
                print(f"[INFO] 查询任务状态 ({attempt + 1}/{max_poll_retries})...")

                query_url = f"{video_api_url}/{task_id}"
                query_response = requests.get(query_url, headers=headers, timeout=poll_timeout)

                if query_response.status_code != 200:
                    logger.warning(f"查询任务状态失败: {query_response.status_code}")
                    # 继续下一次查询，不立即降级
                    time.sleep(poll_interval)
                    continue

                try:
                    query_result = query_response.json()
                    status = query_result.get('status', '').lower()

                    logger.info(f"任务状态: {status}")

                    if status == 'succeeded':
                        video_url = query_result.get('result', {}).get('video_url')
                        if video_url:
                            logger.info(f"下载视频: {video_url}")

                            # 下载视频（带重试）
                            max_download_retries = 3
                            for dl_retry in range(max_download_retries):
                                try:
                                    logger.info(f"下载视频 (尝试 {dl_retry + 1}/{max_download_retries})...")
                                    video_response = requests.get(video_url, timeout=120)

                                    if video_response.status_code == 200:
                                        with open(video_filepath, 'wb') as f:
                                            f.write(video_response.content)
                                        file_size = os.path.getsize(video_filepath)
                                        logger.info(f"视频已保存: {file_size} 字节")
                                        return
                                    else:
                                        logger.warning(f"视频下载失败: {video_response.status_code}")
                                        if dl_retry < max_download_retries - 1:
                                            time.sleep(5)
                                            continue
                                        else:
                                            logger.info("降级到本地视频生成...")
                                            create_intelligent_video(video_filepath, video_data)
                                            return

                                except Exception as e:
                                    logger.warning(f"下载视频时出错: {str(e)}")
                                    if dl_retry < max_download_retries - 1:
                                        time.sleep(5)
                                        continue
                                    else:
                                        logger.info("降级到本地视频生成...")
                                        create_intelligent_video(video_filepath, video_data)
                                        return
                        else:
                            logger.warning("未找到视频URL")
                            logger.info(f"响应内容: {query_result}")
                            break

                    elif status == 'failed':
                        error_info = query_result.get('error', '未知错误')
                        logger.warning(f"任务失败: {error_info}")
                        logger.info(f"响应内容: {query_result}")
                        break

                    elif status in ['pending', 'processing', 'running']:
                        # 任务正在进行中，继续等待
                        time.sleep(poll_interval)

                    else:
                        logger.warning(f"未知的任务状态: {status}")
                        time.sleep(poll_interval)

                except json.JSONDecodeError as e:
                    logger.warning(f"解析查询响应JSON失败: {str(e)}")
                    logger.info(f"响应内容: {query_response.text[:200]}")
                    time.sleep(poll_interval)

            except requests.exceptions.Timeout as e:
                logger.warning(f"查询任务状态超时: {str(e)}")
                time.sleep(poll_interval)

            except Exception as e:
                logger.error(f"查询任务状态时发生错误: {str(e)}")
                import traceback
                traceback.print_exc()
                time.sleep(poll_interval)

        logger.info("轮询超时或任务失败，降级到本地视频生成...")
        create_intelligent_video(video_filepath, video_data)
    else:
        logger.info("未能创建任务，降级到本地视频生成...")
        create_intelligent_video(video_filepath, video_data)

def create_intelligent_video(video_filepath, video_data):
    """创建智能视频（降级方案）- 生成基于文档内容的演示视频"""
    logger.info("使用本地视频生成引擎创建演示视频...")

    try:
        import cv2
        import numpy as np
        from PIL import Image, ImageDraw, ImageFont

        # 解析视频数据
        ai_prompt = video_data.get('ai_prompt', '')
        original_content = video_data.get('original_content', '')
        style = video_data.get('style', 'business')
        duration = video_data.get('duration', 'medium')

        # 时长映射
        duration_map = {'short': 30, 'medium': 60, 'long': 100}
        video_duration = duration_map.get(duration, 60)
        fps = 30
        total_frames = video_duration * fps

        # 视频尺寸
        width, height = 1920, 1080

        # 风格配色
        style_colors = {
            'business': {'bg': (44, 62, 80), 'text': (236, 240, 241), 'accent': (52, 152, 219)},
            'creative': {'bg': (231, 76, 60), 'text': (236, 240, 241), 'accent': (241, 196, 15)},
            'minimalist': {'bg': (255, 255, 255), 'text': (44, 62, 80), 'accent': (52, 152, 219)},
            'tech': {'bg': (44, 62, 80), 'text': (52, 152, 219), 'accent': (46, 204, 113)}
        }
        colors = style_colors.get(style, style_colors['business'])

        # 提取关键内容
        content_lines = [line.strip() for line in original_content.split('\n') if line.strip()][:10]
        document_title = content_lines[0] if content_lines else "文档介绍"
        key_points = content_lines[1:5] if len(content_lines) > 1 else ["文档内容介绍"]

        # 创建视频写入器
        fourcc = cv2.VideoWriter_fourcc(*'mp4v')
        video_writer = cv2.VideoWriter(video_filepath, fourcc, fps, (width, height))

        # 场景配置
        scenes = [
            {'title': document_title, 'content': '', 'duration': 0.2},  # 开场20%
            {'title': '核心内容', 'content': key_points[0] if len(key_points) > 0 else '', 'duration': 0.25},
            {'title': '关键要点', 'content': key_points[1] if len(key_points) > 1 else '', 'duration': 0.25},
            {'title': '重要信息', 'content': key_points[2] if len(key_points) > 2 else '', 'duration': 0.2},
            {'title': '谢谢观看', 'content': '', 'duration': 0.1}  # 结尾10%
        ]

        # 为每个场景生成帧
        current_frame = 0
        for scene in scenes:
            scene_frames = int(total_frames * scene['duration'])

            for i in range(scene_frames):
                # 创建帧
                img = Image.new('RGB', (width, height), colors['bg'])
                draw = ImageDraw.Draw(img)

                # 添加标题
                try:
                    title_font = ImageFont.truetype("arial.ttf", 80)
                    content_font = ImageFont.truetype("arial.ttf", 50)
                except:
                    title_font = ImageFont.load_default()
                    content_font = ImageFont.load_default()

                # 标题位置
                title_bbox = draw.textbbox((0, 0), scene['title'], font=title_font)
                title_width = title_bbox[2] - title_bbox[0]
                title_x = (width - title_width) // 2
                title_y = height // 3

                # 绘制标题
                draw.text((title_x, title_y), scene['title'], fill=colors['text'], font=title_font)

                # 添加内容
                if scene['content']:
                    content_lines_wrap = wrap_text(scene['content'], 30)
                    content_y = title_y + 150

                    for line in content_lines_wrap[:3]:  # 最多显示3行
                        content_bbox = draw.textbbox((0, 0), line, font=content_font)
                        content_width = content_bbox[2] - content_bbox[0]
                        content_x = (width - content_width) // 2
                        draw.text((content_x, content_y), line, fill=colors['accent'], font=content_font)
                        content_y += 70

                # 添加装饰元素
                draw.rectangle([50, 50, width-50, height-50], outline=colors['accent'], width=5)

                # 转换为OpenCV格式
                frame = cv2.cvtColor(np.array(img), cv2.COLOR_RGB2BGR)
                video_writer.write(frame)
                current_frame += 1

        # 释放视频写入器
        video_writer.release()
        logger.info(f"视频生成成功: {video_filepath}")

    except ImportError as e:
        logger.warning(f"缺少必要的库: {str(e)}")
        logger.info("请安装opencv-python和Pillow库: pip install opencv-python Pillow")
        # 创建占位视频
        create_placeholder_video(video_filepath, video_data)
    except Exception as e:
        logger.error(f"视频生成失败: {str(e)}")
        import traceback
        traceback.print_exc()
        create_placeholder_video(video_filepath, video_data)

def wrap_text(text, max_length):
    """文本换行"""
    words = list(text)
    lines = []
    current_line = ''

    for char in words:
        if len(current_line) < max_length:
            current_line += char
        else:
            lines.append(current_line)
            current_line = char

    if current_line:
        lines.append(current_line)

    return lines

def create_placeholder_video(video_filepath, video_data):
    """创建占位视频"""
    logger.info("创建占位视频...")

    try:
        import cv2
        import numpy as np

        width, height = 1280, 720
        fps = 30
        duration = 5  # 5秒占位视频
        total_frames = duration * fps

        fourcc = cv2.VideoWriter_fourcc(*'mp4v')
        video_writer = cv2.VideoWriter(video_filepath, fourcc, fps, (width, height))

        for i in range(total_frames):
            # 创建渐变背景
            frame = np.zeros((height, width, 3), dtype=np.uint8)
            frame[:, :, 0] = int(255 * (i / total_frames))  # R
            frame[:, :, 1] = int(200 * (1 - i / total_frames))  # G
            frame[:, :, 2] = 150  # B

            # 添加文字
            text = "视频生成中..."
            cv2.putText(frame, text, (width//2 - 200, height//2), 
                       cv2.FONT_HERSHEY_SIMPLEX, 2, (255, 255, 255), 3)

            video_writer.write(frame)

        video_writer.release()
        logger.info(f"占位视频创建成功: {video_filepath}")

    except Exception as e:
        logger.error(f"创建占位视频失败: {str(e)}")

# Vercel 需要的导出
app = app

if __name__ == '__main__':
    if os.environ.get('FLASK_ENV') != 'production':
        app.run(debug=True, host='0.0.0.0', port=5000)
    else:
        app.run()
